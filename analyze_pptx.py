#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Скрипт для анализа содержимого PPTX файла
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def analyze_pptx(file_path: str) -> None:
    """Анализирует содержимое PPTX файла и выводит детальную информацию"""
    
    if not Path(file_path).exists():
        print(f"❌ Файл {file_path} не найден!")
        return
    
    try:
        prs = Presentation(file_path)
        print(f"📊 АНАЛИЗ ПРЕЗЕНТАЦИИ: {Path(file_path).name}")
        print("=" * 60)
        
        # Общая информация
        print(f"📋 Общее количество слайдов: {len(prs.slides)}")
        print(f"📐 Размер слайда: {prs.slide_width} x {prs.slide_height} EMU")
        print(f"📏 Размер в дюймах: {prs.slide_width/914400:.2f}\" x {prs.slide_height/914400:.2f}\"")
        print()
        
        # Анализ каждого слайда
        for i, slide in enumerate(prs.slides, 1):
            print(f"🎯 СЛАЙД {i}:")
            print(f"   Layout: {slide.slide_layout.name}")
            print(f"   Количество объектов: {len(slide.shapes)}")
            
            # Анализ заголовка
            title_found = False
            if hasattr(slide.shapes, 'title') and slide.shapes.title:
                if slide.shapes.title.has_text_frame and slide.shapes.title.text_frame.text.strip():
                    print(f"   📝 Заголовок: '{slide.shapes.title.text_frame.text.strip()}'")
                    title_found = True
            
            if not title_found:
                print(f"   ⚠️  Заголовок отсутствует или пустой")
            
            # Анализ содержимого
            text_shapes = 0
            image_shapes = 0
            other_shapes = 0
            total_text_length = 0
            
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_shapes += 1
                    text_content = shape.text_frame.text.strip()
                    total_text_length += len(text_content)
                    if text_content and shape != slide.shapes.title:
                        print(f"   📄 Текст ({len(text_content)} символов): '{text_content[:100]}{'...' if len(text_content) > 100 else ''}'")
                        
                        # Анализ параграфов и буллетов
                        for j, paragraph in enumerate(shape.text_frame.paragraphs):
                            if paragraph.text.strip():
                                level = paragraph.level
                                bullet_info = "с буллетом" if paragraph.level > 0 or (hasattr(paragraph, '_element') and paragraph._element.get('marL')) else "без буллета"
                                print(f"     • Параграф {j+1} (уровень {level}, {bullet_info}): '{paragraph.text.strip()[:50]}{'...' if len(paragraph.text.strip()) > 50 else ''}'")
                
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_shapes += 1
                    print(f"   🖼️  Изображение: {shape.width/914400:.2f}\" x {shape.height/914400:.2f}\"")
                else:
                    other_shapes += 1
            
            print(f"   📊 Статистика: {text_shapes} текстовых блоков, {image_shapes} изображений, {other_shapes} других объектов")
            print(f"   📝 Общий объем текста: {total_text_length} символов")
            
            # Проверка заметок
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                print(f"   📋 Заметки ({len(notes_text)} символов): '{notes_text[:100]}{'...' if len(notes_text) > 100 else ''}'")
            else:
                print(f"   ⚠️  Заметки отсутствуют")
            
            print()
        
        # Общая оценка качества
        print("🔍 ОЦЕНКА КАЧЕСТВА:")
        print("=" * 30)
        
        issues = []
        
        if len(prs.slides) == 0:
            issues.append("❌ Презентация пустая")
        
        empty_slides = 0
        slides_without_title = 0
        slides_without_content = 0
        
        for slide in prs.slides:
            # Проверка заголовка
            has_title = False
            if hasattr(slide.shapes, 'title') and slide.shapes.title:
                if slide.shapes.title.has_text_frame and slide.shapes.title.text_frame.text.strip():
                    has_title = True
            
            if not has_title:
                slides_without_title += 1
            
            # Проверка содержимого
            has_content = False
            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title:
                    if shape.text_frame.text.strip():
                        has_content = True
                        break
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    has_content = True
                    break
            
            if not has_content:
                slides_without_content += 1
            
            if not has_title and not has_content:
                empty_slides += 1
        
        if empty_slides > 0:
            issues.append(f"❌ {empty_slides} полностью пустых слайдов")
        
        if slides_without_title > 0:
            issues.append(f"⚠️  {slides_without_title} слайдов без заголовка")
        
        if slides_without_content > 0:
            issues.append(f"⚠️  {slides_without_content} слайдов без содержимого")
        
        if not issues:
            print("✅ Презентация выглядит корректно")
        else:
            print("Обнаружены проблемы:")
            for issue in issues:
                print(f"   {issue}")
        
        print()
        
    except Exception as e:
        print(f"❌ Ошибка при анализе файла: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Использование: python analyze_pptx.py <путь_к_файлу.pptx>")
        sys.exit(1)
    
    analyze_pptx(sys.argv[1])