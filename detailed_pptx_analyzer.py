#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Расширенный анализатор качества PPTX презентаций
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import re

def analyze_presentation_quality(pptx_path: str):
    """Детальный анализ качества презентации"""
    
    if not Path(pptx_path).exists():
        print(f"❌ Файл не найден: {pptx_path}")
        return
    
    try:
        prs = Presentation(pptx_path)
        
        print(f"📊 ДЕТАЛЬНЫЙ АНАЛИЗ ПРЕЗЕНТАЦИИ: {Path(pptx_path).name}")
        print("=" * 70)
        
        # Общая информация
        slide_count = len(prs.slides)
        print(f"📄 Общее количество слайдов: {slide_count}")
        print(f"📐 Размер слайдов: {prs.slide_width} x {prs.slide_height} EMU")
        print()
        
        # Анализ каждого слайда
        issues = []
        total_text_length = 0
        total_images = 0
        slides_with_titles = 0
        slides_with_content = 0
        
        for i, slide in enumerate(prs.slides, 1):
            print(f"📋 СЛАЙД {i}:")
            print("-" * 40)
            
            # Анализ заголовка
            title_text = ""
            if slide.shapes.title:
                title_text = slide.shapes.title.text.strip()
                if title_text:
                    slides_with_titles += 1
                    print(f"   📝 Заголовок: '{title_text[:60]}{'...' if len(title_text) > 60 else ''}'")
                else:
                    print(f"   ⚠️  Заголовок пустой")
                    issues.append(f"Слайд {i}: пустой заголовок")
            else:
                print(f"   ❌ Заголовок отсутствует")
                issues.append(f"Слайд {i}: отсутствует заголовок")
            
            # Анализ содержимого
            text_shapes = 0
            image_shapes = 0
            other_shapes = 0
            slide_text_length = 0
            
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or hasattr(shape, 'text_frame'):
                    if hasattr(shape, 'text') and shape.text.strip():
                        text_shapes += 1
                        slide_text_length += len(shape.text)
                        
                        # Показать первые строки текста
                        text_preview = shape.text.strip()[:100]
                        if '\n' in text_preview:
                            text_preview = text_preview.split('\n')[0] + '...'
                        print(f"   📄 Текст: '{text_preview}{'...' if len(shape.text) > 100 else ''}'")
                        
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_shapes += 1
                    print(f"   🖼️  Изображение найдено")
                else:
                    other_shapes += 1
            
            total_text_length += slide_text_length
            total_images += image_shapes
            
            if text_shapes > 0 or image_shapes > 0:
                slides_with_content += 1
            
            print(f"   📊 Объекты: {text_shapes} текстовых, {image_shapes} изображений, {other_shapes} других")
            print(f"   📝 Длина текста: {slide_text_length} символов")
            
            # Проверка на проблемы
            if text_shapes == 0 and image_shapes == 0:
                print(f"   ⚠️  Слайд пустой (нет содержимого)")
                issues.append(f"Слайд {i}: пустой слайд")
            
            if slide_text_length > 1000:
                print(f"   ⚠️  Слишком много текста ({slide_text_length} символов)")
                issues.append(f"Слайд {i}: слишком много текста")
            
            # Проверка на дублированный контент
            if '**Слайд' in title_text:
                print(f"   ⚠️  Заголовок содержит метаинформацию")
                issues.append(f"Слайд {i}: заголовок содержит метаинформацию")
            
            print()
        
        # Общая статистика
        print("📈 ОБЩАЯ СТАТИСТИКА:")
        print("=" * 40)
        print(f"📄 Слайдов с заголовками: {slides_with_titles}/{slide_count} ({slides_with_titles/slide_count*100:.1f}%)")
        print(f"📋 Слайдов с содержимым: {slides_with_content}/{slide_count} ({slides_with_content/slide_count*100:.1f}%)")
        print(f"📝 Общая длина текста: {total_text_length} символов")
        print(f"🖼️  Общее количество изображений: {total_images}")
        print(f"📊 Средняя длина текста на слайд: {total_text_length/slide_count:.1f} символов")
        print()
        
        # Анализ проблем
        print("🔍 АНАЛИЗ ПРОБЛЕМ:")
        print("=" * 40)
        
        if not issues:
            print("✅ Серьезных проблем не обнаружено")
        else:
            print(f"⚠️  Найдено проблем: {len(issues)}")
            for issue in issues:
                print(f"   • {issue}")
        
        print()
        
        # Оценка качества
        print("🎯 ОЦЕНКА КАЧЕСТВА:")
        print("=" * 40)
        
        quality_score = 100
        
        # Штрафы за проблемы
        if slides_with_titles < slide_count * 0.8:
            quality_score -= 20
            print("❌ Многие слайды без заголовков (-20 баллов)")
        
        if slides_with_content < slide_count * 0.9:
            quality_score -= 15
            print("❌ Есть пустые слайды (-15 баллов)")
        
        if total_text_length < 100:
            quality_score -= 25
            print("❌ Слишком мало текста (-25 баллов)")
        
        if len(issues) > slide_count * 0.3:
            quality_score -= 20
            print("❌ Много проблем с форматированием (-20 баллов)")
        
        if total_images == 0 and slide_count > 3:
            quality_score -= 10
            print("⚠️  Нет изображений (-10 баллов)")
        
        # Бонусы
        if slides_with_titles == slide_count:
            quality_score += 5
            print("✅ Все слайды имеют заголовки (+5 баллов)")
        
        if total_images > 0:
            quality_score += 5
            print("✅ Есть изображения (+5 баллов)")
        
        quality_score = max(0, min(100, quality_score))
        
        print(f"\n🏆 ИТОГОВАЯ ОЦЕНКА: {quality_score}/100")
        
        if quality_score >= 80:
            print("✅ Отличное качество презентации")
        elif quality_score >= 60:
            print("👍 Хорошее качество презентации")
        elif quality_score >= 40:
            print("⚠️  Удовлетворительное качество, есть проблемы")
        else:
            print("❌ Низкое качество, требуется доработка")
        
        print()
        
        # Рекомендации
        print("💡 РЕКОМЕНДАЦИИ:")
        print("=" * 40)
        
        if slides_with_titles < slide_count:
            print("• Добавьте заголовки ко всем слайдам")
        
        if slides_with_content < slide_count:
            print("• Удалите пустые слайды или добавьте к ним содержимое")
        
        if total_text_length < 500:
            print("• Добавьте больше содержательного текста")
        
        if total_images == 0:
            print("• Рассмотрите возможность добавления изображений")
        
        if any('метаинформацию' in issue for issue in issues):
            print("• Очистите заголовки от служебной информации")
        
        if not issues and quality_score >= 80:
            print("• Презентация готова к использованию!")
        
    except Exception as e:
        print(f"❌ Ошибка при анализе: {e}")
        import traceback
        traceback.print_exc()

def compare_presentations(old_path: str, new_path: str):
    """Сравнивает две презентации"""
    
    print("🔄 СРАВНЕНИЕ ПРЕЗЕНТАЦИЙ")
    print("=" * 50)
    
    try:
        old_prs = Presentation(old_path) if Path(old_path).exists() else None
        new_prs = Presentation(new_path) if Path(new_path).exists() else None
        
        if not old_prs:
            print(f"❌ Старая презентация не найдена: {old_path}")
            return
        
        if not new_prs:
            print(f"❌ Новая презентация не найдена: {new_path}")
            return
        
        old_slides = len(old_prs.slides)
        new_slides = len(new_prs.slides)
        
        print(f"📊 Количество слайдов: {old_slides} → {new_slides} ({new_slides - old_slides:+d})")
        
        # Сравнение содержимого
        old_text_total = 0
        new_text_total = 0
        
        for slide in old_prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    old_text_total += len(shape.text)
        
        for slide in new_prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    new_text_total += len(shape.text)
        
        print(f"📝 Общий объем текста: {old_text_total} → {new_text_total} ({new_text_total - old_text_total:+d})")
        
        if new_slides > old_slides:
            print("✅ Улучшение: больше слайдов")
        elif new_slides < old_slides:
            print("⚠️  Уменьшение количества слайдов")
        
        if new_text_total > old_text_total:
            print("✅ Улучшение: больше текстового содержимого")
        
    except Exception as e:
        print(f"❌ Ошибка при сравнении: {e}")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Использование: python detailed_pptx_analyzer.py <файл.pptx> [старый_файл.pptx]")
        sys.exit(1)
    
    pptx_file = sys.argv[1]
    
    if len(sys.argv) >= 3:
        old_file = sys.argv[2]
        compare_presentations(old_file, pptx_file)
        print()
    
    analyze_presentation_quality(pptx_file)