from __future__ import annotations

from pathlib import Path
from typing import List, Optional

from pptx import Presentation
from pptx.util import Inches

from .models import SlideModel, TextBlock, ImageBlock


def build_presentation(slides: List[SlideModel], out_file: Path, template: Optional[str] = None) -> None:
    prs = Presentation(template) if template else Presentation()
    for idx, slide in enumerate(slides):
        layout = prs.slide_layouts[0] if idx == 0 else prs.slide_layouts[1]
        pptx_slide = prs.slides.add_slide(layout)
        if slide.title:
            title_placeholder = pptx_slide.shapes.title
            title_placeholder.text = slide.title
        body = pptx_slide.placeholders[1] if len(pptx_slide.placeholders) > 1 else None
        if body:
            tf = body.text_frame
            for block in slide.blocks:
                if isinstance(block, TextBlock):
                    if block.bullets:
                        for bullet in block.bullets:
                            p = tf.add_paragraph()
                            p.text = bullet
                            p.level = 0
                    else:
                        p = tf.add_paragraph()
                        p.text = block.text
                elif isinstance(block, ImageBlock):
                    pptx_slide.shapes.add_picture(block.src, Inches(1), Inches(2))
        if slide.notes:
            notes_frame = pptx_slide.notes_slide.notes_text_frame
            notes_frame.text = slide.notes
    prs.save(out_file)
